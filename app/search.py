"""
검색 API 라우터

POST /api/search        : 직장 + 조건 → 카드 리스트 (LLM 친구 한 마디 포함)
GET  /api/apt/{seq}/routes : 단지 상세 경로 옵션 (rank 1~N)
"""
import asyncio
import hashlib
import re
import time
from fastapi import APIRouter, Depends, HTTPException, BackgroundTasks
from pydantic import BaseModel, Field, field_validator
from app.db import get_db, connect as db_connect
from app.workplaces import get_or_create
from app.transit import fetch_cells, haversine, cell_center
from app.ai import build_recommendations, build_stats, build_comments, card_key
from app.portable import upsert_sql, year_month_minus, year_minus, greatest
from config import cfg

router = APIRouter()


# ── 요청 스키마 ──────────────────────────────────────────────
_ALLOWED_PYEONG = {'10평미만', '10평대', '20평대', '30평대', '40평대', '50평대+'}


class SearchRequest(BaseModel):
    workplace_address: str = Field(..., min_length=2, max_length=200, description="직장 도로명/지번 주소")
    workplace_address_2: str | None = Field(
        None, min_length=2, max_length=200,
        description="두 번째 직장 주소 (맞벌이용, 선택)",
    )
    max_minutes:  int = Field(60, ge=10, le=60, description="사용자 선택값. 내부적으로 10분 여유 차감")
    max_price:    int = Field(50000, ge=1000, le=2_000_000, description="만원 단위, 예: 50000=5억")
    pyeong_types: list[str] = Field(
        default_factory=lambda: ['10평대', '20평대'], min_length=1, max_length=6,
    )
    min_kaptdaCnt: int | None = Field(
        None, ge=0, le=100_000, description="최소 단지 세대수 (선택). 미지정=100",
    )
    build_year_min: int | None = Field(
        None, ge=1960, le=2030, description="최소 준공연도 (예: 2010 → 2010년 이후 준공 단지만)",
    )
    min_price: int | None = Field(
        None, ge=1000, le=2_000_000,
        description="최소 가격 (만원). 예: 30000=3억. 미지정=하한 없음.",
    )

    @field_validator('workplace_address_2')
    @classmethod
    def _check_workplace_2(cls, v, info):
        if v is not None:
            wp1 = info.data.get('workplace_address', '').strip()
            if v.strip() == wp1:
                raise ValueError('두 직장이 동일합니다. 단일 모드로 검색하세요.')
        return v

    @field_validator('min_price')
    @classmethod
    def _check_min_price(cls, v, info):
        if v is not None:
            max_p = info.data.get('max_price')
            if max_p is not None and v >= max_p:
                raise ValueError('min_price는 max_price보다 작아야 합니다')
        return v

    @field_validator('pyeong_types')
    @classmethod
    def _check_pyeong_types(cls, v: list[str]) -> list[str]:
        invalid = [p for p in v if p not in _ALLOWED_PYEONG]
        if invalid:
            raise ValueError(f'허용되지 않는 평형: {invalid}')
        return v


# 통근시간 여유분 (사용자 선택값에서 차감)
# 10→15: 60분 입력 시 effective 45분, 반경 16.7→15km, 셀 ~19%↓
COMMUTE_BUFFER_MIN = 15

# Vercel Hobby maxDuration 60s 안에 끝내기 위한 ODsay 호출 셀 상한.
# 신규 워크플레이스 첫 검색에서 초과분은 다음 검색에 캐시 채워짐 (자연 점진 완성).
# 4키 × 4동시 × ~1.2s/round × 100ms sleep → ~250셀이 ~25s 안에 처리됨.
MAX_FETCH_CELLS_PER_CALL = 200   # 4키×4동시×1.2s×round → ~200셀이 ~20s, 여유 확보

# Vercel Hobby maxDuration 60s 제약 관련 상수
ODSAY_HARD_TIMEOUT_S = 30   # ODsay 호출 하드컷 (30s) — 이후 DB쿼리·처리에 ~20s 확보
WALL_CLOCK_BUDGET_S  = 50   # 함수 전체 예산. 초과 시 cards 쿼리 스킵, partial 즉시 반환


# ── POST /api/search ─────────────────────────────────────────
@router.post("/search")
async def search(req: SearchRequest, background_tasks: BackgroundTasks, conn=Depends(get_db)):
    import time as _time
    _t0 = _time.monotonic()   # 벽시계 타이머 시작 (Vercel 60s 예산 추적용)

    # ─ 1. workplace 등록 / wp_id 확보 ─
    # get_or_create 안의 Kakao HTTP 호출이 동기 블로킹이므로 to_thread로 오프로드
    wp = await asyncio.to_thread(get_or_create, conn, req.workplace_address)
    if not wp:
        raise HTTPException(400, f'주소 변환 실패: {req.workplace_address}')

    wp_id = wp['wp_id']
    dest_lat, dest_lng = wp['lat'], wp['lng']
    # 통근시간 여유분 차감: 사용자 60분 입력 → 실제 50분 이내 매물만
    effective_max_min = max(req.max_minutes - COMMUTE_BUFFER_MIN, 5)
    radius_km = effective_max_min * 20 / 60

    # ─ 1b. 두 번째 직장 (dual workplace 모드) ─
    wp_2 = None
    if req.workplace_address_2:
        wp_2 = await asyncio.to_thread(get_or_create, conn, req.workplace_address_2)
        if not wp_2:
            raise HTTPException(400, f'두 번째 주소 변환 실패: {req.workplace_address_2}')
        if wp_2['wp_id'] == wp['wp_id']:
            raise HTTPException(422, '두 직장이 동일합니다. 단일 모드로 검색하세요.')
    dual = wp_2 is not None
    wp_id_2 = wp_2['wp_id'] if dual else None

    # ─ 2. 후보 단지 → 셀 ─
    apt_filter = 'WHERE is_apt=1 AND recent_trade=3'
    apt_params = []
    if req.min_kaptdaCnt is not None:
        apt_filter += ' AND kaptdaCnt >= ?'
        apt_params.append(req.min_kaptdaCnt)
    if req.build_year_min is not None:
        apt_filter += ' AND build_year >= ?'
        apt_params.append(req.build_year_min)
    apts = conn.execute(
        f'SELECT apt_seq, lat, lng FROM apartments {apt_filter}', apt_params
    ).fetchall()
    near_keys = [
        a['apt_seq'] for a in apts
        if haversine(a['lat'], a['lng'], dest_lat, dest_lng) <= radius_km
        and (not dual or haversine(a['lat'], a['lng'], wp_2['lat'], wp_2['lng']) <= radius_km)
    ]

    if not near_keys:
        return _empty_response(wp, [], wp_2=wp_2)

    ph = ','.join('?'*len(near_keys))
    pt = ','.join('?'*len(req.pyeong_types))
    min_price_clause = ' AND deal_amount_int>=?' if req.min_price is not None else ''
    matched_params = near_keys + req.pyeong_types + [req.max_price]
    if req.min_price is not None:
        matched_params.append(req.min_price)
    matched_rows = conn.execute(
        f'SELECT DISTINCT apt_seq FROM trade_recent '
        f'WHERE apt_seq IN ({ph}) AND pyeong_type IN ({pt}) AND deal_amount_int<=?'
        f'{min_price_clause}',
        matched_params
    ).fetchall()
    matched_keys = [r['apt_seq'] for r in matched_rows]

    if not matched_keys:
        return _empty_response(wp, [], wp_2=wp_2)

    ph2 = ','.join('?'*len(matched_keys))
    cells = sorted(set(r['grid_key'] for r in conn.execute(
        f'SELECT DISTINCT grid_key FROM apartments WHERE apt_seq IN ({ph2})',
        matched_keys
    ).fetchall()))

    # ─ 3. 캐시 미스 셀만 ODsay 호출 ─
    # passed_filter=1 인 셀만 "캐시됨"으로 간주 → passed_filter=0/응답실패는 재호출
    cached_1 = set(r['origin_cell'] for r in conn.execute(
        'SELECT origin_cell FROM transit_cache WHERE wp_id=? AND passed_filter=1', (wp_id,)
    ).fetchall())
    conn.execute('DELETE FROM transit_cache WHERE wp_id=? AND passed_filter=0', (wp_id,))
    if dual:
        cached_2 = set(r['origin_cell'] for r in conn.execute(
            'SELECT origin_cell FROM transit_cache WHERE wp_id=? AND passed_filter=1', (wp_id_2,)
        ).fetchall())
        conn.execute('DELETE FROM transit_cache WHERE wp_id=? AND passed_filter=0', (wp_id_2,))
    else:
        cached_2 = set()
    conn.commit()

    # 셀별 직장 거리 계산 헬퍼
    def _cell_dist_to(lat, lng):
        def _fn(c):
            clat, clng = cell_center(c)
            return haversine(clat, clng, lat, lng)
        return _fn

    to_fetch_all_1 = [c for c in cells if c not in cached_1]
    to_fetch_all_2 = [c for c in cells if c not in cached_2] if dual else []

    # ── 동적 분배: wp별 미스 수를 측정 후 적은 쪽이 먼저 할당량 차지 ──
    TOTAL_LIMIT = MAX_FETCH_CELLS_PER_CALL  # 200
    half = TOTAL_LIMIT // 2                  # 100
    n1, n2 = len(to_fetch_all_1), len(to_fetch_all_2)

    if not dual:
        take_1, take_2 = min(n1, TOTAL_LIMIT), 0
    elif n1 <= half and n2 <= half:
        take_1, take_2 = n1, n2
    elif n1 <= half:
        take_1 = n1
        take_2 = min(n2, TOTAL_LIMIT - n1)
    elif n2 <= half:
        take_2 = n2
        take_1 = min(n1, TOTAL_LIMIT - n2)
    else:
        take_1, take_2 = half, half

    to_fetch_1 = sorted(to_fetch_all_1, key=_cell_dist_to(dest_lat, dest_lng))[:take_1]
    to_fetch_2 = (
        sorted(to_fetch_all_2, key=_cell_dist_to(wp_2['lat'], wp_2['lng']))[:take_2]
        if dual else []
    )
    deferred_cells_1 = n1 - len(to_fetch_1)
    deferred_cells_2 = n2 - len(to_fetch_2) if dual else 0

    # 하드 타임아웃: Vercel 60s 제약 안에 반드시 응답 반환.
    try:
        if dual:
            odsay_results = await asyncio.wait_for(
                asyncio.gather(
                    fetch_cells(conn, wp, to_fetch_1),
                    fetch_cells(conn, wp_2, to_fetch_2),
                ),
                timeout=ODSAY_HARD_TIMEOUT_S,
            )
            odsay_stats_1, odsay_stats_2 = odsay_results
        else:
            odsay_stats_1 = await asyncio.wait_for(
                fetch_cells(conn, wp, to_fetch_1),
                timeout=ODSAY_HARD_TIMEOUT_S,
            )
            odsay_stats_2 = None
    except asyncio.TimeoutError:
        print(f'[search] ODsay {ODSAY_HARD_TIMEOUT_S}s 타임아웃 — 캐시 결과로만 응답')
        odsay_stats_1 = {'fetched': 0, 'passed': 0, 'failed': 0,
                         'elapsed_ms': ODSAY_HARD_TIMEOUT_S * 1000}
        odsay_stats_2 = None
        deferred_cells_1 += len(to_fetch_1)
        deferred_cells_2 += len(to_fetch_2)

    # 기존 변수명 호환 (단일 모드 하위 호환)
    odsay_stats = odsay_stats_1
    deferred_cells = deferred_cells_1 + deferred_cells_2

    # ── 벽시계 잔여 예산 체크 ────────────────────────────────────
    # ODsay 이후 cards 쿼리·추천 로직에 최소 10s 필요.
    # 예산 초과 시 cards 쿼리 없이 즉시 partial 반환 → 504 방지.
    _elapsed = _time.monotonic() - _t0
    if _elapsed > WALL_CLOCK_BUDGET_S - 10:
        print(f'[search] 벽시계 예산 초과 ({_elapsed:.1f}s) — cards 쿼리 생략, partial 반환')
        return {
            'wp_id': wp_id, 'llm_pending': False,
            'workplace': {'address_input': wp['address_input'],
                          'address_norm': wp['address_norm'],
                          'lat': wp['lat'], 'lng': wp['lng']},
            'stats': {'total': 0}, 'buckets': [], 'cards': [],
            'meta': {'partial': True,
                     'deferred_cells': deferred_cells + len(cells),
                     'total_cells': len(cells),
                     'odsay_elapsed_ms': int(_elapsed * 1000)},
        }

    # ─ 4. 카드 쿼리 ─
    min_cnt_clause        = ' AND a.kaptdaCnt >= ?' if req.min_kaptdaCnt is not None else ''
    build_year_clause     = ' AND a.build_year >= ?' if req.build_year_min is not None else ''
    min_price_card_clause = ' AND t.deal_amount_int >= ?' if req.min_price is not None else ''

    def _extra_params():
        p = []
        if req.min_kaptdaCnt is not None:
            p.append(req.min_kaptdaCnt)
        if req.build_year_min is not None:
            p.append(req.build_year_min)
        if req.min_price is not None:
            p.append(req.min_price)
        return p

    if not dual:
        # ── 단일 모드 (기존 쿼리 그대로) ────────────────────────
        cards_params = [wp_id, effective_max_min, req.max_price, *req.pyeong_types,
                        *_extra_params()]
        cards = conn.execute(f"""
            SELECT
                a.apt_seq, a.apt_nm, a.umd_nm, a.kaptdaCnt, a.lat, a.lng,
                a.kaptCode,
                r.total_time_min, r.bus_cnt, r.subway_cnt,
                r.step1_type, r.step1_time_min, r.step1_노선,
                r.step2_type, r.step2_time_min, r.step2_노선,
                r.step3_type, r.step3_time_min, r.step3_노선,
                r.step4_type, r.step4_time_min, r.step4_노선,
                r.step5_type, r.step5_time_min, r.step5_노선,
                t.pyeong_type,
                MIN(t.deal_amount_int) AS price_low,
                MAX(t.deal_amount_int) AS price_high,
                COUNT(t.id) AS deal_count,
                AVG(t.deal_amount_int * 1.0 / NULLIF(t.pyeong, 0)) AS pyeong_price_avg,
                MAX(k.kaptTopFloor) AS top_floor,
                MAX(k.kaptUsedate) AS use_date
            FROM apartments a
            LEFT JOIN kapt_complexes k ON a.kaptCode = k.kaptCode
            JOIN transit_routes r ON a.grid_key = r.origin_cell
            JOIN trade_recent t   ON a.apt_seq  = t.apt_seq
            WHERE a.is_apt=1
              AND r.wp_id=? AND r.rank=1 AND r.total_time_min<=?
              AND t.deal_amount_int<=?
              AND t.pyeong_type IN ({pt})
              AND t.dealing_gbn = '중개거래'
              {min_cnt_clause}
              {build_year_clause}
              {min_price_card_clause}
            GROUP BY
                a.apt_seq, a.apt_nm, a.umd_nm, a.kaptdaCnt, a.lat, a.lng,
                a.kaptCode,
                r.total_time_min, r.bus_cnt, r.subway_cnt,
                r.step1_type, r.step1_time_min, r.step1_노선,
                r.step2_type, r.step2_time_min, r.step2_노선,
                r.step3_type, r.step3_time_min, r.step3_노선,
                r.step4_type, r.step4_time_min, r.step4_노선,
                r.step5_type, r.step5_time_min, r.step5_노선,
                t.pyeong_type
            ORDER BY r.total_time_min, price_low
        """, cards_params).fetchall()
    else:
        # ── dual 모드 — r1(wp1) + r2(wp2) self-join ────────────
        # ORDER BY: GREATEST(t1, t2) 를 portable.greatest()로 처리
        order_expr = greatest('r1.total_time_min', 'r2.total_time_min')
        cards_params = [
            wp_id, effective_max_min,
            wp_id_2, effective_max_min,
            req.max_price, *req.pyeong_types,
            *_extra_params(),
        ]
        cards = conn.execute(f"""
            SELECT
                a.apt_seq, a.apt_nm, a.umd_nm, a.kaptdaCnt, a.lat, a.lng,
                a.kaptCode,
                r1.total_time_min AS total_time_min_1,
                r1.bus_cnt AS bus_cnt_1, r1.subway_cnt AS subway_cnt_1,
                r1.step1_type, r1.step1_time_min, r1.step1_노선,
                r1.step2_type, r1.step2_time_min, r1.step2_노선,
                r1.step3_type, r1.step3_time_min, r1.step3_노선,
                r1.step4_type, r1.step4_time_min, r1.step4_노선,
                r1.step5_type, r1.step5_time_min, r1.step5_노선,
                r2.total_time_min AS total_time_min_2,
                r2.bus_cnt AS bus_cnt_2, r2.subway_cnt AS subway_cnt_2,
                r2.step1_type AS step1_type_2, r2.step1_time_min AS step1_time_min_2,
                r2.step1_노선 AS step1_노선_2,
                r2.step2_type AS step2_type_2, r2.step2_time_min AS step2_time_min_2,
                r2.step2_노선 AS step2_노선_2,
                r2.step3_type AS step3_type_2, r2.step3_time_min AS step3_time_min_2,
                r2.step3_노선 AS step3_노선_2,
                r2.step4_type AS step4_type_2, r2.step4_time_min AS step4_time_min_2,
                r2.step4_노선 AS step4_노선_2,
                r2.step5_type AS step5_type_2, r2.step5_time_min AS step5_time_min_2,
                r2.step5_노선 AS step5_노선_2,
                t.pyeong_type,
                MIN(t.deal_amount_int) AS price_low,
                MAX(t.deal_amount_int) AS price_high,
                COUNT(t.id) AS deal_count,
                AVG(t.deal_amount_int * 1.0 / NULLIF(t.pyeong, 0)) AS pyeong_price_avg,
                MAX(k.kaptTopFloor) AS top_floor,
                MAX(k.kaptUsedate) AS use_date
            FROM apartments a
            LEFT JOIN kapt_complexes k ON a.kaptCode = k.kaptCode
            JOIN transit_routes r1 ON a.grid_key = r1.origin_cell
            JOIN transit_routes r2 ON a.grid_key = r2.origin_cell
            JOIN trade_recent t    ON a.apt_seq  = t.apt_seq
            WHERE a.is_apt=1
              AND r1.wp_id=? AND r1.rank=1 AND r1.total_time_min<=?
              AND r2.wp_id=? AND r2.rank=1 AND r2.total_time_min<=?
              AND t.deal_amount_int<=?
              AND t.pyeong_type IN ({pt})
              AND t.dealing_gbn = '중개거래'
              {min_cnt_clause}
              {build_year_clause}
              {min_price_card_clause}
            GROUP BY
                a.apt_seq, a.apt_nm, a.umd_nm, a.kaptdaCnt, a.lat, a.lng,
                a.kaptCode,
                r1.total_time_min, r1.bus_cnt, r1.subway_cnt,
                r1.step1_type, r1.step1_time_min, r1.step1_노선,
                r1.step2_type, r1.step2_time_min, r1.step2_노선,
                r1.step3_type, r1.step3_time_min, r1.step3_노선,
                r1.step4_type, r1.step4_time_min, r1.step4_노선,
                r1.step5_type, r1.step5_time_min, r1.step5_노선,
                r2.total_time_min, r2.bus_cnt, r2.subway_cnt,
                r2.step1_type, r2.step1_time_min, r2.step1_노선,
                r2.step2_type, r2.step2_time_min, r2.step2_노선,
                r2.step3_type, r2.step3_time_min, r2.step3_노선,
                r2.step4_type, r2.step4_time_min, r2.step4_노선,
                r2.step5_type, r2.step5_time_min, r2.step5_노선,
                t.pyeong_type
            ORDER BY {order_expr}, price_low
        """, cards_params).fetchall()

    # ─ 4b. 카드별 최근 거래 (3개월, 최대 4건) ─
    # 카드 단위 = (apt_seq, pyeong_type) 이므로 recent_map 키도 동일하게.
    if cards:
        seqs = list({c['apt_seq'] for c in cards})
        ph = ','.join('?' * len(seqs))
        threshold_ym = year_month_minus(3)  # 3개월 전 YYYY*100+MM
        # ROW_NUMBER()로 DB에서 최대 4건만 가져옴 (전체 fetchall 후 Python 필터 방식 대비
        # 대형 단지/활발한 지역에서 수만 행 전송을 방지)
        recent_rows = conn.execute(f"""
            WITH ranked AS (
                SELECT apt_seq, pyeong, pyeong_type, floor, deal_amount_int,
                       deal_year, deal_month, deal_day, dealing_gbn,
                       ROW_NUMBER() OVER (
                           PARTITION BY apt_seq, pyeong_type
                           ORDER BY deal_year DESC, deal_month DESC, deal_day DESC
                       ) AS rn
                FROM trade_recent
                WHERE apt_seq IN ({ph})
                  AND deal_year*100 + deal_month >= ?
                  AND dealing_gbn = '중개거래'
            )
            SELECT apt_seq, pyeong, pyeong_type, floor, deal_amount_int,
                   deal_year, deal_month, deal_day, dealing_gbn
            FROM ranked WHERE rn <= 4
            ORDER BY apt_seq, pyeong_type,
                     deal_year DESC, deal_month DESC, deal_day DESC
        """, seqs + [threshold_ym]).fetchall()
        recent_map: dict = {}
        for row in recent_rows:
            key = (row['apt_seq'], row['pyeong_type'])
            recent_map.setdefault(key, []).append({
                'pyeong': row['pyeong'],
                'floor': row['floor'],
                'amount': row['deal_amount_int'],
                'date': (
                    f"{row['deal_year'] % 100:02d}"
                    f".{row['deal_month']:02d}.{row['deal_day']:02d}"
                ),
                'gbn': row['dealing_gbn'],
            })
    else:
        recent_map = {}

    # ─ 4b-2. 3개월 평균가 맵 (배지 대표 가격용) ─
    avg_price_map: dict = {
        key: round(sum(t['amount'] for t in trades) / len(trades))
        for key, trades in recent_map.items() if trades
    }

    # ─ 4c. why_tags 배치 조회 ─
    # trade_tags 테이블이 없거나 비어 있으면 graceful degradation → why_tags: []
    tag_map: dict = {}
    if cards:
        seqs = list({c['apt_seq'] for c in cards})
        ph = ','.join('?' * len(seqs))
        try:
            tag_rows = conn.execute(
                f'SELECT apt_seq, pyeong_type, tag_type, label, detail '
                f'FROM trade_tags WHERE apt_seq IN ({ph})',
                seqs,
            ).fetchall()
            for row in tag_rows:
                key = (row['apt_seq'], row['pyeong_type'])
                tag_map.setdefault(key, []).append({
                    'type':   row['tag_type'],
                    'label':  row['label'],
                    'detail': row['detail'],
                })
        except Exception:
            # trade_tags 미존재 시 빈 tag_map 유지.
            # ⚠️ Postgres(pgBouncer) 에서는 쿼리 실패 후 트랜잭션이 aborted 상태로 남는다.
            # rollback() 없이 다음 쿼리를 시도하면 InFailedSqlTransaction 500 이 발생하므로
            # 여기서 반드시 롤백하여 커넥션 상태를 복구한다.
            try:
                conn.rollback()
            except Exception:
                pass

    # ─ 4d. 가격 변동률 배치 조회 (trade_history 최근 3개월 vs 4~9개월 전) ─
    price_chg_map: dict = {}
    if cards:
        import datetime
        _now = datetime.date.today()

        def _ym(months_back: int) -> int:
            """months_back 개월 전 YYYYMM 정수 반환."""
            y, m = _now.year, _now.month - months_back
            while m <= 0:
                m += 12
                y -= 1
            return y * 100 + m

        ym_3mo = _ym(3)   # 3개월 전 (최근 구간 시작)
        ym_9mo = _ym(9)   # 9개월 전 (과거 구간 시작)
        seqs_chg = list({c['apt_seq'] for c in cards})
        ph_chg = ','.join('?' * len(seqs_chg))
        try:
            chg_rows = conn.execute(
                f'SELECT apt_seq, pyeong_type,'
                f'  AVG(CASE WHEN deal_year*100+deal_month >= {ym_3mo}'
                f'           THEN deal_amount_int END) AS recent_avg,'
                f'  AVG(CASE WHEN deal_year*100+deal_month >= {ym_9mo}'
                f'            AND deal_year*100+deal_month < {ym_3mo}'
                f'           THEN deal_amount_int END) AS past_avg'
                f' FROM trade_history'
                f' WHERE apt_seq IN ({ph_chg})'
                f'   AND deal_year*100+deal_month >= {ym_9mo}'
                f' GROUP BY apt_seq, pyeong_type',
                seqs_chg,
            ).fetchall()
            for row in chg_rows:
                r_avg, p_avg = row['recent_avg'], row['past_avg']
                if r_avg and p_avg and p_avg > 0:
                    pct = round((r_avg - p_avg) / p_avg * 100, 1)
                    if abs(pct) >= 3.0:
                        price_chg_map[(row['apt_seq'], row['pyeong_type'])] = pct
        except Exception:
            try:
                conn.rollback()
            except Exception:
                pass

    # ─ 5. 카드 변환 + 추천 로직 (통근버킷 × 평형 매트릭스) ─
    raw_cards = [_card_to_dict(c, recent_map, tag_map, price_chg_map, avg_price_map, dual=dual) for c in cards]
    # 대표가(3개월 평균)가 검색 범위를 벗어난 카드 제거
    # avg_price_map이 없는 카드(거래 없음)는 유지
    def _price_in_range(card: dict) -> bool:
        p = card.get('price_low')
        if not p:
            return True
        # avg_price_map에서 온 값만 필터 (recent_map 기반으로 계산된 경우)
        key = (card['apt_seq'], card['pyeong_type'])
        if key not in avg_price_map:
            return True  # fallback 값이면 필터 안 함
        if req.max_price and p > req.max_price:
            return False
        if req.min_price and p < req.min_price:
            return False
        return True
    raw_cards = [c for c in raw_cards if _price_in_range(c)]
    rec = build_recommendations(raw_cards, effective_max_min)
    buckets = rec['buckets']
    all_cards = rec['cards']

    # ─ 6. 통계 계산 ─
    stats = build_stats(all_cards, buckets)

    # ─ 7. LLM 친구 한 마디 — 전체 카드 (추천=Sonnet 긴 코멘트 / 일반=Haiku 한 줄) ─
    cached_comments: dict[str, str] = {}  # card_key → comment
    if all_cards:
        # wp_id 단위로 모든 캐시 조회 (카드 수 많아도 단일 쿼리)
        rows = conn.execute(
            'SELECT apt_seq, pyeong_type, comment FROM apt_pt_friend_comment '
            'WHERE wp_id=?',
            [wp_id],
        ).fetchall()
        for row in rows:
            cached_comments[f"{row['apt_seq']}:{row['pyeong_type']}"] = row['comment']

    miss_cards = [c for c in all_cards if card_key(c) not in cached_comments]
    llm_pending = len(miss_cards) > 0
    if miss_cards:
        wp_label = wp.get('address_norm') or wp.get('address_input', '')
        background_tasks.add_task(_generate_comments_bg, miss_cards, all_cards, wp_id, wp_label)

    # 카드에 코멘트 병합 (캐시된 것만 즉시 표시)
    for c in all_cards:
        c['friend_comment'] = cached_comments.get(card_key(c), '')

    return {
        'wp_id': wp_id,
        'llm_pending': llm_pending,
        'workplace': {
            'address_input': wp['address_input'],
            'address_norm':  wp['address_norm'],
            'lat': wp['lat'], 'lng': wp['lng'],
        },
        'workplace_2': {
            'wp_id': wp_2['wp_id'],
            'address_input': wp_2['address_input'],
            'address_norm':  wp_2['address_norm'],
            'lat': wp_2['lat'], 'lng': wp_2['lng'],
        } if dual else None,
        'stats': stats,
        'buckets': buckets,
        'cards': all_cards,
        'meta': {
            'dual_workplace':   dual,
            'odsay_calls_made': odsay_stats['fetched'],
            'odsay_passed':     odsay_stats['passed'],
            'odsay_failed':     odsay_stats['failed'],
            'odsay_calls_made_1': odsay_stats_1['fetched'],
            'odsay_calls_made_2': odsay_stats_2['fetched'] if odsay_stats_2 else 0,
            'cache_hits':         len(cached_1) + len(cached_2),
            'odsay_elapsed_ms':   odsay_stats['elapsed_ms'],
            'llm_cached':       len(cached_comments),
            'llm_pending':      len(miss_cards),
            'llm_pending_recommend': sum(1 for c in miss_cards if c.get('is_recommended')),
            'llm_pending_regular':   sum(1 for c in miss_cards if not c.get('is_recommended')),
            'partial':          deferred_cells > 0,
            'deferred_cells':   deferred_cells,
            'deferred_cells_1': deferred_cells_1,
            'deferred_cells_2': deferred_cells_2,
            'total_cells':      len(cells),
        },
    }


# ── 백그라운드: LLM 코멘트 생성 + DB 캐싱 ───────────────────
async def _generate_comments_bg(miss_cards: list, all_cards: list, wp_id: int, wp_label: str):
    """BackgroundTask — 응답 후 실행. 별도 DB 커넥션 사용.

    miss_cards: 코멘트 생성 대상 (추천 카드 중 캐시 미스)
    all_cards : 전체 카드 (평형별 평균가 등 컨텍스트 계산용)
    """
    import time
    import traceback
    t0 = time.time()
    rec_n = sum(1 for c in miss_cards if c.get('is_recommended'))
    reg_n = len(miss_cards) - rec_n
    print(f'[bg_comments] 시작 — 추천 {rec_n}개, 일반 {reg_n}개')
    try:
        new_comments = await build_comments(miss_cards, all_cards, wp_label)
        ok_n = sum(1 for v in new_comments.values()
                   if v.get('comment') and not v['comment'].startswith('('))
        print(f'[bg_comments] LLM 완료 ({time.time()-t0:.1f}s): {ok_n}/{len(new_comments)} 성공')
    except Exception as e:
        print(f'[bg_comments] LLM 호출 자체 실패: {type(e).__name__}: {e}')
        print(traceback.format_exc())
        return

    # DB 저장 (실패 코멘트도 저장하면 다음에 또 시도 안 함 → 성공한 것만 저장)
    try:
        conn = db_connect()
        try:
            rows = []
            for c in miss_cards:
                ck = card_key(c)
                comment = new_comments.get(ck, {}).get('comment', '')
                # '(생성 실패)' 같은 에러 메시지는 캐시하지 않음 → 다음 검색 때 재시도
                if comment and not comment.startswith('('):
                    rows.append((c['apt_seq'], c['pyeong_type'], wp_id, comment))
            if rows:
                conn.executemany(
                    upsert_sql(
                        'apt_pt_friend_comment',
                        ['apt_seq', 'pyeong_type', 'wp_id', 'comment'],
                        pk_cols=['apt_seq', 'pyeong_type', 'wp_id'],
                    ),
                    rows,
                )
                conn.commit()
                print(f'[bg_comments] DB 저장 완료 — {len(rows)}건')
        finally:
            conn.close()
    except Exception as e:
        print(f'[bg_comments] DB 저장 실패: {type(e).__name__}: {e}')
        print(traceback.format_exc())


# ── GET /api/comments — 생성된 댓글 폴링용 ──────────────────
@router.get("/comments")
def get_comments(wp_id: int, keys: str, conn=Depends(get_db)):
    """LLM 백그라운드 완료 후 프론트가 폴링하는 경량 엔드포인트.

    keys: "apt_seq:pyeong_type,apt_seq:pyeong_type,..." 형식
    반환: {"apt_seq:pyeong_type": {"comment": "..."}}
    """
    pairs = []
    for k in keys.split(','):
        k = k.strip()
        if ':' not in k:
            continue
        seq, pt = k.split(':', 1)
        pairs.append((seq.strip(), pt.strip()))
    if not pairs:
        return {}
    if len(pairs) > 200:
        raise HTTPException(400, 'keys 개수가 200을 초과할 수 없습니다')
    conds = ' OR '.join(['(apt_seq=? AND pyeong_type=?)'] * len(pairs))
    params: list = [wp_id]
    for s, p in pairs:
        params.extend([s, p])
    rows = conn.execute(
        f'SELECT apt_seq, pyeong_type, comment FROM apt_pt_friend_comment '
        f"WHERE wp_id=? AND ({conds}) AND comment != ''",
        params,
    ).fetchall()
    return {
        f"{r['apt_seq']}:{r['pyeong_type']}": {'comment': r['comment']}
        for r in rows
    }


def _empty_response(wp, _cards, wp_2=None):
    dual = wp_2 is not None
    return {
        'wp_id': wp['wp_id'],
        'llm_pending': False,
        'workplace': {
            'address_input': wp['address_input'],
            'address_norm':  wp['address_norm'],
            'lat': wp['lat'], 'lng': wp['lng'],
        },
        'workplace_2': {
            'wp_id': wp_2['wp_id'],
            'address_input': wp_2['address_input'],
            'address_norm':  wp_2['address_norm'],
            'lat': wp_2['lat'], 'lng': wp_2['lng'],
        } if dual else None,
        'stats': {'total': 0},
        'buckets': [],
        'cards': [],
        'meta': {'dual_workplace': dual},
    }


def _build_transit_summary(steps: list[dict], bc: int, sc: int, total_time_min: int) -> str:
    """steps 리스트 → 대중교통 요약 문자열 (wp1/wp2 공통 헬퍼)."""
    walk_min = None
    subway_line = None
    for s in steps:
        if s['type'] in ('도보', '환승도보', 'WALK') and walk_min is None:
            walk_min = s['time']
        if s['type'] in ('지하철', 'SUBWAY') and subway_line is None:
            line_raw = s['line'] or ''
            m = re.search(r'\d+호선', line_raw)
            subway_line = m.group(0) if m else (line_raw.strip() or None)

    if subway_line:
        xfer = max(sc - 1, 0)
        summary = f'{subway_line} 직통' if xfer == 0 else f'{subway_line} {xfer}회환승'
        if walk_min:
            summary += f' (도보 {walk_min}분)'
    elif bc:
        xfer = max(bc - 1, 0)
        summary = '버스 직통' if xfer == 0 else f'버스 {xfer}회환승'
        if walk_min:
            summary += f' (도보 {walk_min}분)'
    else:
        summary = f'{total_time_min}분'
    return summary


def _card_to_dict(r, recent_map: dict | None = None, tag_map: dict | None = None, price_chg_map: dict | None = None, avg_price_map: dict | None = None, dual: bool = False):
    # ── wp1 transit 파싱 ─────────────────────────────────────────
    if dual:
        bc_1, sc_1 = r['bus_cnt_1'], r['subway_cnt_1']
        steps_1 = []
        for i in range(1, 6):
            st = r[f'step{i}_type']
            if not st:
                break
            steps_1.append({'type': st, 'time': r[f'step{i}_time_min'], 'line': r[f'step{i}_노선']})
        transit_summary_1 = _build_transit_summary(steps_1, bc_1, sc_1, r['total_time_min_1'])

        # ── wp2 transit 파싱 ────────────────────────────────────
        bc_2, sc_2 = r['bus_cnt_2'], r['subway_cnt_2']
        steps_2 = []
        for i in range(1, 6):
            st = r[f'step{i}_type_2']
            if not st:
                break
            steps_2.append({'type': st, 'time': r[f'step{i}_time_min_2'], 'line': r[f'step{i}_노선_2']})
        transit_summary_2 = _build_transit_summary(steps_2, bc_2, sc_2, r['total_time_min_2'])

        total_time_min = max(r['total_time_min_1'], r['total_time_min_2'])
        # 단일 모드 호환 필드 (transit_summary = wp1 요약)
        bc, sc = bc_1, sc_1
        transit_summary = transit_summary_1
    else:
        bc, sc = r['bus_cnt'], r['subway_cnt']
        steps = []
        for i in range(1, 6):
            st = r[f'step{i}_type']
            if not st:
                break
            steps.append({'type': st, 'time': r[f'step{i}_time_min'], 'line': r[f'step{i}_노선']})
        transit_summary = _build_transit_summary(steps, bc, sc, r['total_time_min'])
        total_time_min = r['total_time_min']
        bc_1 = bc; sc_1 = sc
        transit_summary_1 = transit_summary
        bc_2 = None; sc_2 = None
        transit_summary_2 = None

    # 준공연도: kaptUsedate = "20040517" → 2004
    use_date = r['use_date'] or ''
    build_year = int(use_date[:4]) if len(use_date) >= 4 and use_date[:4].isdigit() else None

    apt_seq = r['apt_seq']
    pyeong_type = r['pyeong_type'] if 'pyeong_type' in r.keys() else None
    try:
        pa = r['pyeong_price_avg']
        pyeong_price = int(pa) if pa else None
    except (IndexError, KeyError):
        pyeong_price = None

    card = {
        'apt_seq': apt_seq, 'apt_nm': r['apt_nm'], 'umd_nm': r['umd_nm'],
        'pyeong_type': pyeong_type,
        'kaptdaCnt': int(r['kaptdaCnt'] or 0),
        'top_floor': r['top_floor'],
        'build_year': build_year,
        'lat': r['lat'], 'lng': r['lng'],
        'total_time_min': total_time_min,
        'total_time_min_1': r['total_time_min_1'] if dual else total_time_min,
        'total_time_min_2': r['total_time_min_2'] if dual else None,
        'bus_cnt': bc, 'subway_cnt': sc,
        'bus_cnt_1': bc_1, 'subway_cnt_1': sc_1,
        'bus_cnt_2': bc_2, 'subway_cnt_2': sc_2,
        'transit_summary': transit_summary,
        'transit_summary_1': transit_summary_1,
        'transit_summary_2': transit_summary_2,
        # 대표가: 3개월 평균 → 최신 거래 → MIN(fallback)
        'price_low': (avg_price_map or {}).get((apt_seq, pyeong_type))
                     or ((recent_map or {}).get((apt_seq, pyeong_type)) or [{}])[0].get('amount')
                     or r['price_low'],
        'price_high': r['price_high'],
        'pyeong_price_avg': pyeong_price,
        'deal_count': r['deal_count'],
        'recent_trades': (recent_map or {}).get((apt_seq, pyeong_type), []),
        'why_tags': (tag_map or {}).get((apt_seq, pyeong_type), []),
        'price_chg_6m_pct': (price_chg_map or {}).get((apt_seq, pyeong_type)),
    }
    return card


# ── GET /api/apt/{apt_seq}/routes ────────────────────────────
@router.get("/apt/{apt_seq}/routes")
def apt_routes(apt_seq: str, wp_id: int, conn=Depends(get_db)):
    """단지 상세 — 모든 경로 옵션 (rank 1~N)"""
    # 뷰(v_apt_transit_options)는 SQLite 쿼리 플래너가 apt_seq 조건을 뷰 안으로 밀어넣지 못해
    # 전체 스캔 → 7초 소요. 2단계 조회(grid_key 먼저 → transit_routes 직접)로 0.001초로 단축.
    gk_row = conn.execute(
        "SELECT grid_key FROM apartments WHERE apt_seq=?", [apt_seq]
    ).fetchone()
    if not gk_row or not gk_row['grid_key']:
        return {'apt_seq': apt_seq, 'wp_id': wp_id, 'options': []}

    rows = conn.execute("""
        SELECT rank, total_time_min, bus_cnt, subway_cnt,
               step1_type, step1_time_min, step1_dist_m, step1_노선, step1_출발, step1_도착,
               step2_type, step2_time_min, step2_dist_m, step2_노선, step2_출발, step2_도착,
               step3_type, step3_time_min, step3_dist_m, step3_노선, step3_출발, step3_도착,
               step4_type, step4_time_min, step4_dist_m, step4_노선, step4_출발, step4_도착,
               step5_type, step5_time_min, step5_dist_m, step5_노선, step5_출발, step5_도착
        FROM transit_routes
        WHERE origin_cell=? AND wp_id=?
        ORDER BY rank
    """, [gk_row['grid_key'], wp_id]).fetchall()

    options = []
    for r in rows:
        steps = []
        for i in range(5):
            off = 4 + i*6   # 각 step당 6개 컬럼 (type, time_min, dist_m, 노선, 출발, 도착)
            t = r[off]
            if not t:
                continue
            steps.append({
                'type': t,
                'time_min': r[off+1],
                'dist_m':   r[off+2],
                'line':     r[off+3],
                'from':     r[off+4],
                'to':       r[off+5],
            })
        options.append({
            'rank': r['rank'],
            'total_time_min': r['total_time_min'],
            'bus_cnt': r['bus_cnt'],
            'subway_cnt': r['subway_cnt'],
            'steps': steps,
        })
    return {'apt_seq': apt_seq, 'wp_id': wp_id, 'options': options}


# ── GET /api/search/apt-lookup  (spec-26: 분석결과 직접 검색) ──
@router.get("/search/apt-lookup")
def search_apt_lookup(
    name: str = "",
    wp_id: int = 0,
    pyeong_type: str = "20평대",
    conn=Depends(get_db),
):
    if len(name.strip()) < 2:
        return {"results": []}

    apts = conn.execute(
        "SELECT apt_seq, apt_nm, umd_nm, lat, lng, kaptdaCnt, build_year "
        "FROM apartments WHERE apt_nm LIKE ? "
        "ORDER BY kaptdaCnt DESC LIMIT 5",
        [f"%{name.strip()}%"],
    ).fetchall()

    results = []
    for apt in apts:
        trade = conn.execute(
            "SELECT price_low, price_high FROM trade_recent "
            "WHERE apt_seq=? AND pyeong_type=? LIMIT 1",
            [apt["apt_seq"], pyeong_type],
        ).fetchone()

        transit = conn.execute(
            "SELECT MIN(total_time) AS transit_min FROM transit_cache "
            "WHERE wp_id=? AND passed_filter=1 "
            "  AND origin_cell=(SELECT grid_key FROM apartments WHERE apt_seq=? LIMIT 1)",
            [wp_id, apt["apt_seq"]],
        ).fetchone()

        results.append(
            {
                "apt_seq":    apt["apt_seq"],
                "apt_nm":     apt["apt_nm"],
                "umd_nm":     apt["umd_nm"],
                "lat":        apt["lat"],
                "lng":        apt["lng"],
                "kaptdaCnt":  apt["kaptdaCnt"],
                "build_year": apt["build_year"],
                "pyeong_type": pyeong_type,
                "price_low":  trade["price_low"]  if trade else None,
                "price_high": trade["price_high"] if trade else None,
                "transit_min": transit["transit_min"] if transit and transit["transit_min"] else None,
            }
        )

    return {"results": results}


# ── GET /api/apt/{apt_seq}/detail ────────────────────────────
@router.get("/apt/{apt_seq}/detail")
def apt_detail(apt_seq: str, wp_id: int, conn=Depends(get_db)):
    """
    상세 패널용 — 거래내역 + POI + 시세차트 데이터
    통근경로는 /routes 엔드포인트를 별도 호출
    """
    import time
    timings = {}
    t0 = time.time()

    # ── 기본 단지 정보 + 건물 상세 ────────────────────────────────
    apt = conn.execute("""
        SELECT a.apt_nm, a.umd_nm, a.kaptdaCnt, a.lat, a.lng,
               k.kaptUsedate, k.kaptTopFloor, k.kaptBaseFloor,
               k.kaptDongCnt, k.kaptdEcnt,
               k.kaptdCccnt, k.kaptdPcnt, k.kaptdPcntu,
               k.codeHeatNm, k.codeHallNm,
               k.kaptBcompany,
               k.groundElChargerCnt, k.undergroundElChargerCnt,
               k.subwayLine, k.subwayStation, k.kaptdWtimesub
        FROM apartments a
        LEFT JOIN kapt_complexes k USING(kaptCode)
        WHERE a.apt_seq = ?
    """, [apt_seq]).fetchone()
    timings['apt_info'] = round((time.time()-t0)*1000)
    if not apt:
        print(f'[detail {apt_seq}] 단지 없음 — timings: {timings}')
        return {}

    # 주차대수: kaptdPcnt(지하) + kaptdPcntu(지상)
    def _to_int(v):
        try:
            return int(v or 0)
        except (ValueError, TypeError):
            return 0
    parking_underground = _to_int(apt['kaptdCccnt'])
    parking_ground = _to_int(apt['kaptdPcntu'])
    parking_total = (
        parking_underground + parking_ground
        if (parking_underground or parking_ground) else None
    )

    # 준공연도
    use_date = apt['kaptUsedate'] or ''
    build_year = int(use_date[:4]) if len(use_date) >= 4 and use_date[:4].isdigit() else None

    # 전기차 충전기
    ev_chargers = _to_int(apt['groundElChargerCnt']) + _to_int(apt['undergroundElChargerCnt'])

    building_info = {
        'build_year':   build_year,
        'top_floor':    apt['kaptTopFloor'],
        'base_floor':   apt['kaptBaseFloor'],
        'dong_cnt':     apt['kaptDongCnt'],
        'elevator_cnt': apt['kaptdEcnt'],
        'parking':      parking_total,
        'heat_type':    apt['codeHeatNm'],
        'hall_type':    apt['codeHallNm'],
        'builder':      apt['kaptBcompany'],
        'ev_chargers':  ev_chargers if ev_chargers else None,
        'subway_line':  apt['subwayLine'],
        'subway_sta':   apt['subwayStation'],
        'subway_walk':  apt['kaptdWtimesub'],
    }

    umd_nm = apt['umd_nm']

    # ── 친구 한 마디 (새 테이블 = 평형별 분리) ──────────────────
    # detail 화면은 단지 통합이라 가장 긴 코멘트(=추천 카드의 Sonnet) 우선 표시
    fc_rows = conn.execute(
        'SELECT pyeong_type, comment FROM apt_pt_friend_comment '
        'WHERE apt_seq=? AND wp_id=? ORDER BY LENGTH(comment) DESC',
        [apt_seq, wp_id]
    ).fetchall()
    friend_comment = fc_rows[0]['comment'] if fc_rows else None
    tier = None  # 신규 컨셉엔 tier 없음
    timings['comment'] = round((time.time()-t0)*1000)

    # ── 평수 탭 목록 (거래 많은 순, 디폴트=1위) ─────────────────
    t1 = time.time()
    threshold_year = year_minus(3)  # 3년 전 YYYY
    pyeong_tabs = conn.execute("""
        SELECT pyeong_type, pyeong,
               COUNT(*) AS deal_count,
               MIN(deal_amount_int) AS price_min,
               MAX(deal_amount_int) AS price_max
        FROM trade_history
        WHERE apt_seq = ?
          AND deal_year >= ?
        GROUP BY pyeong_type, pyeong
        ORDER BY deal_count DESC
    """, [apt_seq, threshold_year]).fetchall()

    tabs = [
        {
            'pyeong_type': r['pyeong_type'],
            'pyeong': r['pyeong'],
            'deal_count': r['deal_count'],
            'price_min': r['price_min'],
            'price_max': r['price_max'],
        }
        for r in pyeong_tabs
    ]
    timings['pyeong_tabs'] = round((time.time()-t1)*1000)
    t1 = time.time()

    # ── 시세 차트 — 단지 월별 평균 (3년치, 평수별 전체금액) ──────
    # ym 문자열은 SELECT 시 deal_year/deal_month만 가져와서 Python에서 포맷
    # (printf는 SQLite 전용, Postgres는 LPAD — portable subset에서 빠짐)
    chart_rows = conn.execute("""
        SELECT pyeong_type, pyeong, deal_year, deal_month,
               ROUND(AVG(deal_amount_int)) AS avg_amount,
               COUNT(*) AS cnt
        FROM trade_history
        WHERE apt_seq = ?
          AND deal_year >= ?
        GROUP BY pyeong_type, pyeong, deal_year, deal_month
        ORDER BY pyeong_type, deal_year, deal_month
    """, [apt_seq, threshold_year]).fetchall()

    # 동 평균 — trade_history.umd_nm 직접 사용 (JOIN 없이 idx_th_umd_py_ym 활용)
    dong_avg_rows = conn.execute("""
        SELECT pyeong_type, deal_year, deal_month,
               ROUND(AVG(deal_amount_int)) AS avg_amount,
               COUNT(*) AS cnt
        FROM trade_history
        WHERE umd_nm = ?
          AND deal_year >= ?
        GROUP BY pyeong_type, deal_year, deal_month
        ORDER BY pyeong_type, deal_year, deal_month
    """, [umd_nm, threshold_year]).fetchall()
    timings['chart_rows+dong_avg'] = round((time.time()-t1)*1000)
    t1 = time.time()

    # dict 구조로 변환 {pyeong_type: {ym: avg_amount}}
    complex_chart: dict = {}
    for r in chart_rows:
        key = f"{r['pyeong_type']}_{r['pyeong']}"
        ym = f"{r['deal_year']}-{r['deal_month']:02d}"
        complex_chart.setdefault(key, {})[ym] = int(r['avg_amount'])

    dong_chart: dict = {}
    for r in dong_avg_rows:
        ym = f"{r['deal_year']}-{r['deal_month']:02d}"
        dong_chart.setdefault(r['pyeong_type'], {})[ym] = int(r['avg_amount'])

    # ── 최근 실거래 내역 (최근 20건) ──────────────────────────
    # trade_recent(최신 데이터) UNION trade_history(3년치) → 중복 제거 후 최신순
    trade_rows = conn.execute("""
        SELECT deal_year, deal_month, deal_day,
               pyeong, pyeong_type, deal_amount_int, floor
        FROM trade_recent
        WHERE apt_seq = ?
        UNION
        SELECT deal_year, deal_month, deal_day,
               pyeong, pyeong_type, deal_amount_int, floor
        FROM trade_history
        WHERE apt_seq = ?
          AND deal_year >= ?
        ORDER BY deal_year DESC, deal_month DESC, deal_day DESC
        LIMIT 20
    """, [apt_seq, apt_seq, threshold_year]).fetchall()

    trades = [
        {
            'date': f"{r['deal_year']}.{r['deal_month']:02d}.{r['deal_day']:02d}",
            'pyeong': r['pyeong'],
            'pyeong_type': r['pyeong_type'],
            'amount': r['deal_amount_int'],
            'floor': r['floor'],
        }
        for r in trade_rows
    ]
    timings['trades'] = round((time.time()-t1)*1000)
    t1 = time.time()

    # ── 도보 POI ───────────────────────────────────────────────
    poi_rows = conn.execute("""
        SELECT poi_lclas_cd, poi_mlsfc_cd, poi_nm, distance_m, walking_min
        FROM apt_walking_poi
        WHERE kaptCode = (
            SELECT kaptCode FROM apartments WHERE apt_seq = ? LIMIT 1
        )
          AND walking_min <= 10
        ORDER BY distance_m
        LIMIT 50
    """, [apt_seq]).fetchall()

    poi = [
        {
            'category': r['poi_lclas_cd'],
            'sub_category': r['poi_mlsfc_cd'],
            'name': r['poi_nm'],
            'distance_m': r['distance_m'],
            'walking_min': r['walking_min'],
        }
        for r in poi_rows
    ]

    # ── 시세 요약 지표 (6개월 변화, 동 대비) ──────────────────
    # 가장 많이 거래된 평수로 계산
    price_summary = {}
    if tabs:
        top_pyeong_type = tabs[0]['pyeong_type']
        top_pyeong      = tabs[0]['pyeong']
        key = f"{top_pyeong_type}_{top_pyeong}"
        series = complex_chart.get(key, {})
        months_sorted = sorted(series.keys())
        if len(months_sorted) >= 7:
            recent  = series[months_sorted[-1]]
            six_ago = series[months_sorted[-7]]
            chg_pct = round((recent - six_ago) / six_ago * 100, 1) if six_ago else None
        else:
            chg_pct = None

        dong_series  = dong_chart.get(top_pyeong_type, {})
        dong_months  = sorted(dong_series.keys())
        dong_recent  = dong_series.get(dong_months[-1]) if dong_months else None
        apt_recent   = series.get(months_sorted[-1]) if months_sorted else None
        vs_dong_pct = (
            round((apt_recent - dong_recent) / dong_recent * 100, 1)
            if (apt_recent and dong_recent) else None
        )
        vs_dong_diff = (
            int(apt_recent - dong_recent)
            if (apt_recent and dong_recent) else None
        )

        price_summary = {
            'pyeong_type': top_pyeong_type,
            'pyeong': top_pyeong,
            'change_6m_pct': chg_pct,
            'vs_dong_pct':   vs_dong_pct,
            'vs_dong_diff':  vs_dong_diff,
        }

    timings['poi+summary'] = round((time.time()-t1)*1000)
    timings['TOTAL'] = round((time.time()-t0)*1000)
    print(f'[detail {apt_seq}] timings(ms): {timings}')

    return {
        'apt_seq':   apt_seq,
        'apt_nm':    apt['apt_nm'],
        'umd_nm':    umd_nm,
        'kaptdaCnt': apt['kaptdaCnt'],
        'lat':       apt['lat'],
        'lng':       apt['lng'],
        'tier':           tier,
        'friend_comment': friend_comment,
        'building':       building_info,
        'pyeong_tabs':    tabs,
        'chart': {
            'complex': complex_chart,
            'dong':    dong_chart,
        },
        'price_summary': price_summary,
        'trades': trades,
        'poi':    poi,
    }


# ── 채팅 인메모리 캐시 (spec-23) ────────────────────────────────
_chat_cache: dict = {}
_CHAT_TTL = 3600


def _chat_cache_get(key: tuple) -> "str | None":
    item = _chat_cache.get(key)
    if item and (time.time() - item[0]) < _CHAT_TTL:
        return item[1]
    return None


def _chat_cache_set(key: tuple, reply: str) -> None:
    now = time.time()
    expired = [k for k, v in _chat_cache.items() if now - v[0] >= _CHAT_TTL]
    for k in expired:
        del _chat_cache[k]
    _chat_cache[key] = (now, reply)


def _do_search(query: str) -> str:
    """DuckDuckGo Instant Answer API — API 키 불필요 (spec-24 F3)."""
    import urllib.request
    import urllib.parse
    import json as _json
    try:
        q = urllib.parse.quote(query)
        url = f"https://api.duckduckgo.com/?q={q}&format=json&no_html=1&skip_disambig=1"
        req = urllib.request.Request(url, headers={"User-Agent": "badugi-chat/1.0"})
        with urllib.request.urlopen(req, timeout=5) as r:
            data = _json.loads(r.read())
        abstract   = data.get("AbstractText", "").strip()
        source_url  = data.get("AbstractURL", "").strip()
        source_name = data.get("AbstractSource", "").strip()
        related = [
            t.get("Text", "")
            for t in data.get("RelatedTopics", [])[:3]
            if isinstance(t, dict) and t.get("Text")
        ]
        result = abstract or "\n".join(related)
        if result and source_url:
            label = source_name or source_url
            result += f"\n[출처: {label}] {source_url}"
        return result or "관련 공식 정보를 찾지 못했어."
    except Exception as e:
        return f"검색 오류: {type(e).__name__}"


def _extract_doc_text(data: bytes, media_type: str, filename: str) -> str:
    """Word / PPT / PDF / 텍스트 → 최대 3000자 텍스트 추출 (spec-27)."""
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    try:
        if ext == "docx" or "wordprocessingml" in media_type:
            from docx import Document
            import io as _io
            doc = Document(_io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())[:3000]
        if ext == "pptx" or "presentationml" in media_type:
            from pptx import Presentation
            import io as _io
            prs = Presentation(_io.BytesIO(data))
            texts = [
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            return "\n".join(texts)[:3000]
        if ext == "pdf" or media_type == "application/pdf":
            from pypdf import PdfReader
            import io as _io
            reader = PdfReader(_io.BytesIO(data))
            return "\n".join(
                page.extract_text() or "" for page in reader.pages[:10]
            ).strip()[:3000]
        if ext == "txt" or media_type.startswith("text/"):
            return data.decode("utf-8", errors="ignore")[:3000]
    except Exception as e:
        return f"[파일 파싱 오류: {type(e).__name__}]"
    return ""


def _parse_reply(raw: str) -> "tuple[str, list[str]]":
    """CHIPS 줄을 분리해 (reply, suggestions) 반환."""
    chips_match = re.search(r"\nCHIPS:\s*(.+)$", raw, re.MULTILINE)
    if not chips_match:
        return raw.strip(), []
    suggestions = [s.strip() for s in chips_match.group(1).split("|")][:3]
    suggestions = [s for s in suggestions if s]
    return raw[: chips_match.start()].strip(), suggestions


# ── POST /api/apt/{apt_seq}/chat  (spec-22: 친구 채팅) ─────────
class AptChatRequest(BaseModel):
    pyeong_type: str | None = None
    wp_id:       int | None = None
    message:     str = Field(default="", max_length=500)
    history:     list[dict] = Field(default_factory=list)
    attachments: list[dict] | None = Field(default=None)  # spec-27: [{type, media_type, data, filename}]


_SEARCH_TOOLS = [
    {
        "name": "search_web",
        "description": "학군·지역정보·개발계획·최신뉴스 등 실시간 정보 검색 (DuckDuckGo)",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "검색어 (한국어 포함 가능)"}
            },
            "required": ["query"],
        },
    }
]


@router.post("/apt/{apt_seq}/chat")
def apt_chat(apt_seq: str, req: AptChatRequest, conn=Depends(get_db)):
    import anthropic as _anth
    import datetime as _dt
    import os as _os

    apt = conn.execute(
        "SELECT apt_nm, umd_nm, kaptdaCnt, build_year FROM apartments WHERE apt_seq=? LIMIT 1",
        [apt_seq],
    ).fetchone()
    if not apt:
        raise HTTPException(status_code=404, detail="단지를 찾을 수 없어요")

    history_len = len(req.history or [])
    has_attachment = bool(req.attachments)
    cache_key = (apt_seq, hashlib.md5(req.message.encode()).hexdigest(), history_len)
    if not has_attachment and history_len <= 2:
        cached = _chat_cache_get(cache_key)
        if cached:
            reply, suggestions = _parse_reply(cached)
            return {'reply': reply, 'suggestions': suggestions}

    def _fmt(v: int) -> str:
        e, m = v // 10000, v % 10000
        return f"{e}억{f' {m:,}만' if m else ''}"

    today = _dt.date.today()
    threshold_year = today.year - 1
    d6 = today - _dt.timedelta(days=183)
    d3 = today - _dt.timedelta(days=91)

    stat_rows = conn.execute("""
        SELECT pyeong_type, pyeong,
               ROUND(AVG(deal_amount_int)) AS avg_amt,
               MIN(deal_amount_int) AS min_amt,
               MAX(deal_amount_int) AS max_amt,
               COUNT(*) AS cnt
        FROM trade_history
        WHERE apt_seq = ? AND deal_year >= ?
        GROUP BY pyeong_type, pyeong
        ORDER BY cnt DESC
    """, [apt_seq, threshold_year]).fetchall()

    trend_rows = conn.execute("""
        SELECT deal_year, deal_month, pyeong_type,
               ROUND(AVG(deal_amount_int)) AS avg_amt
        FROM trade_history
        WHERE apt_seq = ?
          AND (deal_year > ? OR (deal_year = ? AND deal_month >= ?))
        GROUP BY deal_year, deal_month, pyeong_type
        ORDER BY pyeong_type, deal_year, deal_month
    """, [apt_seq, d6.year, d6.year, d6.month]).fetchall()

    dong_avg_rows = conn.execute("""
        SELECT pyeong_type, ROUND(AVG(deal_amount_int)) AS avg_amt
        FROM trade_history
        WHERE umd_nm = ?
          AND (deal_year > ? OR (deal_year = ? AND deal_month >= ?))
        GROUP BY pyeong_type
    """, [apt['umd_nm'], d6.year, d6.year, d6.month]).fetchall()

    trades = conn.execute(
        "SELECT pyeong_type, pyeong, deal_year, deal_month, deal_day, deal_amount_int, floor "
        "FROM trade_history WHERE apt_seq=? "
        "ORDER BY deal_year DESC, deal_month DESC, deal_day DESC LIMIT 20",
        [apt_seq],
    ).fetchall()

    stat_lines = [
        f"- {r['pyeong_type']}({r['pyeong']:.0f}평): "
        f"평균 {_fmt(int(r['avg_amt']))} / "
        f"최저 {_fmt(int(r['min_amt']))} / "
        f"최고 {_fmt(int(r['max_amt']))} / "
        f"거래 {r['cnt']}건"
        for r in stat_rows
    ]

    trend_by_type: dict = {}
    for r in trend_rows:
        pt = r['pyeong_type']
        ym = (r['deal_year'], r['deal_month'])
        trend_by_type.setdefault(pt, {})[ym] = int(r['avg_amt'])

    change_lines = []
    for pt, monthly in trend_by_type.items():
        recent = [v for (y, m), v in monthly.items()
                  if y > d3.year or (y == d3.year and m >= d3.month)]
        prev = [v for (y, m), v in monthly.items()
                if not (y > d3.year or (y == d3.year and m >= d3.month))]
        if recent and prev:
            chg = round((sum(recent) / len(recent) - sum(prev) / len(prev))
                        / (sum(prev) / len(prev)) * 100, 1)
            sign = '+' if chg > 0 else ''
            change_lines.append(f"- {pt}: {sign}{chg}%")

    dong_lines = [
        f"- {apt['umd_nm']} {r['pyeong_type']}: 평균 {_fmt(int(r['avg_amt']))}"
        for r in dong_avg_rows
    ]

    trade_lines = "\n".join(
        f"- {r['pyeong_type']}({r['pyeong']:.0f}평) "
        f"{r['deal_year']}.{r['deal_month']:02d}.{r['deal_day']:02d} "
        f"{_fmt(r['deal_amount_int'])} {r['floor']}층"
        for r in trades
    ) or "실거래 데이터 없음"

    # 도보 주요 시설 (10분 이내)
    poi_rows = conn.execute("""
        SELECT poi_lclas_cd, poi_nm, walking_min
        FROM apt_walking_poi
        WHERE kaptCode = (
            SELECT kaptCode FROM apartments WHERE apt_seq = ? LIMIT 1
        )
          AND walking_min <= 10
        ORDER BY walking_min, poi_lclas_cd
        LIMIT 20
    """, [apt_seq]).fetchall()
    poi_lines = [
        f"- {r['poi_nm']} ({r['poi_lclas_cd']}, 도보 {r['walking_min']}분)"
        for r in poi_rows
    ]

    system = f"""너는 부동산을 잘 아는 친한 친구야. 아래 아파트 정보를 바탕으로 친구처럼 솔직하게 답해줘.
반말, 카톡 말투. 5줄 이내.
모르는 건 search_web 도구로 검색해서 답해. 검색 결과가 없거나 제한적이면 "공식 확인이 필요해" 라고 표현해 ("검색이 잘 안 나와" 같은 말은 절대 쓰지 마).
실시간 호가·전세 정보는 없으니 추정할 때 반드시 "확인 필요"를 붙여.
교통 호재·재개발·정부 정책을 언급할 땐 반드시 출처를 명시해. 예: [출처: 기사제목 또는 URL]

답변 마지막에 공백 한 줄 후 반드시 이 형식으로 한국어 후속 질문 3개를 추가해:
CHIPS: 질문1 | 질문2 | 질문3

== 단지 ==
{apt['apt_nm']} · {apt['umd_nm']} · {(apt['kaptdaCnt'] or 0):,}세대 · 준공 {apt['build_year'] or '미상'}년

== 시세 통계 (1년 기준) ==
{chr(10).join(stat_lines) or '데이터 없음'}

== 6개월 시세 변동 ==
{chr(10).join(change_lines) or '데이터 부족'}

== 동 평균 시세 (최근 6개월) ==
{chr(10).join(dong_lines) or '데이터 없음'}

== 최근 실거래 (최대 20건) ==
{trade_lines}

== 도보 10분 이내 주요 시설 ==
{chr(10).join(poi_lines) or '시설 데이터 없음'}"""

    # spec-27: 첨부 파일 처리
    import base64 as _b64
    image_blocks: list[dict] = []
    doc_texts: list[str] = []
    for att in (req.attachments or [])[:1]:  # 최대 1개
        att_data = att.get("data", "")
        if not att_data:
            continue
        att_bytes = _b64.b64decode(att_data)
        att_type = att.get("type", "")
        att_media = att.get("media_type", "")
        att_fname = att.get("filename", "파일")
        if att_type == "image":
            image_blocks.append({
                "type": "image",
                "source": {"type": "base64", "media_type": att_media, "data": att_data},
            })
        else:
            text = _extract_doc_text(att_bytes, att_media, att_fname)
            if text:
                doc_texts.append(f"== 첨부 파일 ({att_fname}) ==\n{text}")

    # 문서 텍스트를 시스템 프롬프트에 추가
    if doc_texts:
        system += "\n\n" + "\n\n".join(doc_texts)

    messages: list[dict] = []
    for h in (req.history or []):
        if h.get('role') in ('user', 'assistant') and h.get('content'):
            messages.append({'role': h['role'], 'content': str(h['content'])})

    # 이미지 첨부 시 vision content blocks
    user_text = req.message or "이 파일을 분석해줘."
    if image_blocks:
        messages.append({'role': 'user', 'content': image_blocks + [{"type": "text", "text": user_text}]})
    else:
        messages.append({'role': 'user', 'content': user_text})

    raw = "잠깐, 다시 시도해봐."
    try:
        api_key = _os.environ.get('ANTHROPIC_API_KEY') or cfg.ANTHROPIC_API_KEY
        client = _anth.Anthropic(api_key=api_key)

        # F3: tool_use agentic loop (최대 3턴)
        MAX_TOOL_TURNS = 3
        for _ in range(MAX_TOOL_TURNS + 1):
            msg = client.messages.create(
                model='claude-opus-4-8',
                max_tokens=700,
                system=system,
                messages=messages,
                tools=_SEARCH_TOOLS,
            )
            if msg.stop_reason != 'tool_use':
                break
            # 검색 도구 실행
            messages.append({'role': 'assistant', 'content': msg.content})
            tool_results = []
            for block in msg.content:
                if block.type == 'tool_use':
                    search_result = _do_search(block.input.get('query', ''))
                    tool_results.append({
                        'type': 'tool_result',
                        'tool_use_id': block.id,
                        'content': search_result,
                    })
            messages.append({'role': 'user', 'content': tool_results})

        text_parts = [b.text for b in msg.content if hasattr(b, 'text') and b.text]
        raw = '\n'.join(text_parts).strip() or "잠깐, 다시 시도해봐."
    except Exception as e:
        raw = f"에러가 났어. 잠깐 기다려봐. ({type(e).__name__}: {str(e)[:80]})"

    if not has_attachment and history_len <= 2:
        _chat_cache_set(cache_key, raw)

    reply, suggestions = _parse_reply(raw)
    return {'reply': reply, 'suggestions': suggestions}
